"""
PPTX转换服务

负责处理HTML到PPTX的转换和文件上传
"""

import hashlib
import threading
import traceback
from pathlib import Path
from typing import Any, Dict, List, Optional

from loguru import logger
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt

from app.core.entity.aigc_metadata import AigcMetadataParams
from app.service import convert_task_manager
from app.service.file_convert.base_convert_service import BaseConvertService, ViewportSize


class PptxConvertService(BaseConvertService):
    """PPTX转换服务类"""

    def __init__(self):
        """初始化PPTX转换服务"""
        super().__init__("PPTX")

    @staticmethod
    def _check_batch_files_exist(batch_dir: Path) -> bool:
        """
        检查批次目录中是否存在相关文件（PPTX服务的实现）

        Args:
            batch_dir: 批次目录路径

        Returns:
            bool: 是否存在相关文件
        """
        try:
            # 检查截图目录的文件名匹配
            screenshots_dir = batch_dir / "screenshots"
            if screenshots_dir.exists():
                screenshot_files = [f.stem for f in screenshots_dir.glob("*.png")]
                # 如果有相同的截图文件名，认为是同一组文件
                return len(screenshot_files) > 0

            return False

        except Exception as e:
            logger.debug(f"检查PPTX批次目录时发生错误: {e}")
            return False

    async def convert_file_keys_to_pptx(
        self,
        file_keys: List[Dict[str, str]],
        task_key: Optional[str] = None,
        sts_credential: Optional[Dict[str, Any]] = None,
        aigc_params: Optional[AigcMetadataParams] = None,
    ) -> Dict[str, Any]:
        """
        将文件key列表转换为PPTX

        Args:
            file_keys: 文件key列表，每个元素包含file_key
            task_key: 任务标识符，会在结果中原样返回
            sts_credential: STS临时凭证，用于上传
            aigc_params: AIGC元数据参数对象

        Returns:
            转换结果字典

        Raises:
            ValueError: 输入参数错误
            RuntimeError: 转换过程中发生不可恢复的错误
        """
        # PPTX转换服务直接确定自己的配置
        options = {}  # 不使用任何默认配置，在转换时动态确定
        # PPTX转换策略：固定1920x1080截图 → PPT页面，无边距布局

        # 通用前置处理：STS解析、file_keys校验、AgentContext初始化、批次目录创建
        sts_cred_obj, batch_id, batch_dir = await self._prepare_conversion_context(
            file_keys=file_keys, sts_credential=sts_credential
        )

        try:
            # 解析文件key到workspace路径
            logger.info("开始解析文件key到workspace路径...")
            try:
                file_path_mapping = await self.resolve_file_keys_to_workspace_paths(file_keys)
            except Exception as e:
                logger.error(f"文件key解析过程中发生错误: {e}")
                raise RuntimeError(f"文件key解析失败: {str(e)}")

            if not file_path_mapping:
                raise RuntimeError("没有找到任何有效的workspace文件，请检查文件key是否正确")

            logger.info(f"开始处理 {len(file_path_mapping)} 个文件")

            # 1. 分析文件，检测PPT入口文件和按目录分组（使用基类通用方法）
            ppt_projects = await self.analyze_and_group_files(
                file_path_mapping, supported_extensions=[".html", ".htm"], service_type="PPTX转换"
            )

            logger.info(f"文件分析完成，发现 {len(ppt_projects)} 个PPT项目")
            for project_name, project_info in ppt_projects.items():
                logger.info(f"  项目 '{project_name}': {len(project_info['files'])} 个文件")

            if not ppt_projects:
                raise RuntimeError("没有找到任何可转换的PPT项目")

            # 2. 计算文件统计信息和并发数（使用基类通用方法）
            total_valid_files, file_type_counts, optimal_concurrency = self.calculate_file_statistics_and_concurrency(
                ppt_projects, supported_extensions=[".html", ".htm"], service_type="PPTX转换"
            )

            # 3. 使用抽象方法执行项目转换
            # 使用全局task_manager实例（如果有task_key）
            task_mgr = convert_task_manager.task_manager if task_key else None

            generated_pptx_files, all_conversion_errors = await self._convert_projects(
                ppt_projects,
                batch_dir,
                options,
                task_mgr,
                task_key,
                total_valid_files,
                optimal_concurrency,
                aigc_params,
            )

            logger.info(f"所有PPT项目处理完成: 成功 {len(generated_pptx_files)}/{len(ppt_projects)} 个项目")

            # 如果有转换错误，记录到日志
            if all_conversion_errors:
                logger.warning(f"部分文件转换失败: {len(all_conversion_errors)} 个错误")
                for error in all_conversion_errors:
                    logger.warning(f"  - {error}")

            # 如果没有任何项目转换成功，抛出错误
            if not generated_pptx_files:
                error_message = "所有PPT项目转换都失败了"
                if all_conversion_errors:
                    error_message += f"，错误详情: {'; '.join(all_conversion_errors)}"
                raise RuntimeError(error_message)

            # 3. 使用抽象方法获取服务特定数据
            service_specific_data = await self._get_service_specific_result_data(
                file_keys, ppt_projects, generated_pptx_files
            )

            # 🎯 计算实际成功处理的输入文件数量
            # 如果没有转换错误，所有有效文件都成功处理
            if not all_conversion_errors:
                actual_processed_files_count = total_valid_files
            else:
                # 有转换错误的情况下，计算成功项目包含的文件数量
                successful_files_count = 0
                for project_name, project_info in ppt_projects.items():
                    # 检查这个项目是否有生成对应的PPTX文件
                    project_has_errors = any(project_name in error for error in all_conversion_errors)
                    if not project_has_errors:
                        successful_files_count += len(project_info.get("files", []))
                actual_processed_files_count = successful_files_count

            result = await self._process_conversion_result(
                converted_files=generated_pptx_files,
                batch_id=batch_id,
                batch_dir=batch_dir,
                # 使用实际文件数量，不是file_keys数量
                valid_files_count=total_valid_files,
                task_key=task_key,
                sts_cred_obj=sts_cred_obj,
                service_specific_data=service_specific_data,
                conversion_errors=all_conversion_errors if all_conversion_errors else None,
                actual_processed_files_count=actual_processed_files_count,
            )

            # 在返回结果中加入total_files字段，应该是实际处理的文件总数
            result["total_files"] = total_valid_files  # 实际处理的文件数量，不是输入file_keys数量

            # 🎯 success_count现在由基类自动处理，使用actual_processed_files_count参数

            # 记录详细的转换统计信息（在设置正确的success_count之后记录）
            logger.info(
                f"PPTX转换统计 - 输入Keys: {len(file_keys)}, 实际文件: {total_valid_files}, "
                f"PPT项目: {len(ppt_projects)}, 成功生成PPTX: {len(generated_pptx_files)}, "
                f"成功处理文件: {result['success_count']}, 转换率: {result['conversion_rate']}%"
            )

            if aigc_params and (aigc_params.user_id or aigc_params.organization_code):
                logger.info(f"PPTX元数据设置 - 用户: {aigc_params.user_id}, 组织: {aigc_params.organization_code}")

            return result

        except Exception as e:
            # 记录详细的错误信息
            logger.error(f"PPTX转换过程中发生错误: {e}")
            logger.error(traceback.format_exc())
            raise
        finally:
            # 保留临时文件，不再清理
            logger.info(f"保留临时文件在目录: {batch_dir}")

    async def _create_pptx_for_project(
        self,
        project_name: str,
        project_info: Dict[str, Any],
        screenshot_dir: Path,
        batch_dir: Path,
        task_mgr=None,
        task_key: Optional[str] = None,
        current_progress: int = 0,
        screenshot_concurrency: int = 1,
        total_valid_files: int = 0,
        aigc_params: Optional[AigcMetadataParams] = None,
    ) -> tuple[Optional[Path], List[str]]:
        """
        为单个PPT项目创建PPTX文件

        Args:
            project_name: 项目名称
            project_info: 项目信息
            screenshot_dir: 截图目录
            batch_dir: 批次目录
            task_mgr: 任务管理器
            task_key: 任务键
            current_progress: 当前进度
            screenshot_concurrency: 截图专用并发数
            total_valid_files: 整个任务的总文件数量（用于正确计算进度）
            aigc_params: AIGC元数据参数对象

        Returns:
            (PPTX文件路径, 错误列表)
        """
        # 使用基类方法准备有序的文件列表
        files_to_convert = self.prepare_ordered_files_for_project(project_info)
        conversion_errors = []

        logger.info(f"开始为项目 '{project_name}' 生成截图，共 {len(files_to_convert)} 个文件")

        # 如果是PPT入口项目，记录排序信息
        if project_info["is_ppt_entry"] and project_info["slides"]:
            logger.info(f"PPT入口项目 '{project_name}' 按slides顺序转换: {len(files_to_convert)} 个文件")

        # 生成截图
        screenshot_files, errors = await self._convert_files_to_screenshots_for_project(
            project_name,
            files_to_convert,
            screenshot_dir,
            task_mgr,
            task_key,
            screenshot_concurrency,
            total_valid_files,
            current_progress,
            aigc_params,
        )

        conversion_errors.extend(errors)

        if not screenshot_files:
            logger.error(f"项目 '{project_name}' 没有成功生成任何截图")
            return None, conversion_errors

        # 创建PPTX文件
        try:
            prs = Presentation()
            blank_slide_layout = prs.slide_layouts[6]  # 6号通常是空白布局

            # 标记是否已设置presentation尺寸
            prs_size_set = False

            # 将截图添加到PPTX
            for screenshot_file in screenshot_files:
                slide = prs.slides.add_slide(blank_slide_layout)
                try:
                    img = Image.open(screenshot_file)
                    img_width, img_height = img.size

                    # 如果是第一张图片，根据图片尺寸设置presentation的宽高
                    if not prs_size_set:
                        # 使用常规PPT单位设置presentation尺寸
                        # 标准转换: 96 DPI 下，1像素 = 0.75点 (72/96 = 0.75)

                        # 将像素转换为点 (Pt) 单位
                        pixels_to_pt = 72.0 / 96.0  # 0.75
                        width_pt = img_width * pixels_to_pt
                        height_pt = img_height * pixels_to_pt

                        # 使用Pt单位设置精确尺寸
                        prs.slide_width = Pt(width_pt)
                        prs.slide_height = Pt(height_pt)
                        prs_size_set = True

                        logger.info(
                            f"设置presentation尺寸: {img_width}x{img_height}px → {width_pt:.1f}x{height_pt:.1f}pt"
                        )

                    # 将图片设置得比slide稍大，消除右侧和底部白边
                    # 按图片原始比例等比超出，避免图片变形
                    img_ratio = img_width / img_height

                    # 基础超出量，然后按比例分配
                    base_extra = Pt(6)  # 基础超出量约8px

                    # 按图片宽高比计算等比例的宽度和高度增量
                    if img_ratio >= 1:  # 横图或正方形
                        width_extra = base_extra
                        height_extra = base_extra / img_ratio
                    else:  # 竖图
                        height_extra = base_extra
                        width_extra = base_extra * img_ratio

                    slide.shapes.add_picture(
                        str(screenshot_file),
                        Inches(0),  # 左侧对齐，不超出
                        Inches(0),  # 上侧对齐，不超出
                        width=prs.slide_width + width_extra,  # 按比例增加宽度
                        height=prs.slide_height + height_extra,  # 按比例增加高度
                    )
                except Exception as img_error:
                    logger.warning(f"无法获取图片尺寸，使用默认设置: {img_error}")
                    # 如果还没有设置presentation尺寸，使用默认的16:9宽屏尺寸
                    if not prs_size_set:
                        prs.slide_width = Inches(16.0)
                        prs.slide_height = Inches(9.0)
                        prs_size_set = True
                        logger.info("无法获取图片尺寸，使用默认16:9宽屏设置 (16.0x9.0英寸)")

                    # 回退方法，假设16:9比例进行等比超出
                    assumed_ratio = 16.0 / 9.0  # 假设16:9横图比例
                    base_extra = Pt(6)  # 基础超出量

                    # 按假设比例计算等比例增量
                    width_extra = base_extra
                    height_extra = base_extra / assumed_ratio

                    slide.shapes.add_picture(
                        str(screenshot_file),
                        Inches(0),  # 左侧对齐，不超出
                        Inches(0),  # 上侧对齐，不超出
                        width=prs.slide_width + width_extra,  # 按比例增加宽度
                        height=prs.slide_height + height_extra,  # 按比例增加高度
                    )

            # 保存PPTX文件 - 使用时间戳命名格式
            pptx_filename = self._generate_timestamped_filename(project_name, "pptx")
            pptx_path = batch_dir / pptx_filename
            prs.save(str(pptx_path))
            logger.info(f"项目 '{project_name}' PPTX文件创建成功: {pptx_path}")

            # 嵌入AIGC签名元数据到PPTX文件（必须成功，否则整个转换失败）
            await self.embed_pptx_metadata(str(pptx_path), aigc_params)
            logger.info(f"PPT转换：成功嵌入签名元数据到PPTX文件 {pptx_path}")

            # 不再使用SHA256重命名，保持时间戳文件名
            return pptx_path, conversion_errors

        except Exception as e:
            logger.error(f"为项目 '{project_name}' 创建PPTX文件失败: {e}")
            conversion_errors.append(f"项目 '{project_name}' PPTX创建失败: {str(e)}")
            return None, conversion_errors

    async def _convert_files_to_screenshots_for_project(
        self,
        project_name: str,
        files_to_convert: List[Path],
        screenshot_dir: Path,
        task_mgr=None,
        task_key: Optional[str] = None,
        screenshot_concurrency: int = 1,
        total_valid_files: int = 0,
        progress_offset: int = 0,
        aigc_params: Optional[AigcMetadataParams] = None,
    ) -> tuple[List[Path], List[str]]:
        """
        为单个项目的文件生成截图

        Args:
            project_name: 项目名称
            files_to_convert: 要转换的文件列表
            screenshot_dir: 截图输出目录
            task_mgr: 任务管理器
            task_key: 任务键
            screenshot_concurrency: 截图专用并发数（避免内存爆炸）
            total_valid_files: 整个任务的总文件数量（用于正确计算进度）
            progress_offset: 进度偏移量，用于跨项目累积进度
            aigc_params: AIGC元数据参数对象

        Returns:
            (截图文件列表, 错误列表)
        """
        screenshot_files = []
        conversion_errors = []

        logger.info(f"项目 '{project_name}' 使用截图专用并发，并发数: {screenshot_concurrency} (内存安全)")

        # 用于进度跟踪的共享状态 - 从progress_offset开始计数，支持跨项目累积
        progress_lock = threading.Lock()
        completed_screenshots = progress_offset  # 🎯 从偏移量开始，而不是从0开始

        # 使用共享浏览器进行转换（支持服务端模式）
        playwright_instance = None
        shared_browser = None
        shared_context = None
        try:
            logger.info(f"项目 '{project_name}': 准备共享浏览器实例（支持服务端模式）")

            # 使用类型安全的视口配置
            viewport = ViewportSize(width=1920, height=1080)

            playwright_instance, shared_browser, shared_context = await self._create_shared_browser_context(
                browser_type="screenshot",
                viewport=viewport,
                device_scale_factor=2.0,
                context_options={
                    "color_scheme": "light",
                    "reduced_motion": "no-preference",
                    # 语言和编码设置 - 修复PPT导出乱码
                    "locale": "zh-CN",
                    "timezone_id": "Asia/Shanghai",
                    "extra_http_headers": {
                        "Accept-Language": "zh-CN,zh;q=0.9,en-US;q=0.8,en;q=0.7",
                        "Accept-Charset": "utf-8,*;q=0.5",
                    },
                },
            )

            async def process_single_file_to_screenshot(
                source_path: Path, index: int
            ) -> tuple[Optional[Path], Optional[str]]:
                """处理单个文件的截图转换 - 串行处理避免网络资源竞争"""
                nonlocal completed_screenshots
                try:
                    logger.debug(f"项目 '{project_name}' 串行处理第 {index + 1} 个文件: {source_path}")

                    if not source_path.exists():
                        return None, f"文件不存在: {source_path}"

                    file_suffix = source_path.suffix.lower()
                    if file_suffix not in [".html", ".htm"]:
                        return None, f"跳过非HTML文件: {source_path}"

                    page = await shared_context.new_page()
                    page.set_default_timeout(self.PAGE_OPERATION_TIMEOUT)
                    self._bind_page_console_logger(page, debug_info=f"PPTX截图-{source_path.name}")

                    try:
                        # 🎯 使用基类的增强页面加载方法，包含智能外部资源和字体加载优化
                        debug_info = f"PPTX截图{index + 1}"
                        success = await self._load_html_page_standard(page, source_path, debug_info)

                        if not success:
                            return None, f"页面加载失败: {source_path}"

                        # 保持网页原始CSS样式，不进行任何修改
                        logger.debug(f"PPT截图 {index + 1}: 页面加载和渲染完全完成（包含外部资源优化）")

                        # 生成截图
                        path_hash = hashlib.md5(str(source_path).encode()).hexdigest()[:8]
                        screenshot_filename = f"{project_name}_{source_path.stem}_{path_hash}.png"
                        screenshot_output_path = screenshot_dir / screenshot_filename

                        # 使用viewport截图确保精确的1920x1080尺寸，避免white边
                        await page.screenshot(
                            path=str(screenshot_output_path),
                            full_page=False,  # 使用viewport截图，精确控制尺寸
                            timeout=self.SCREENSHOT_TIMEOUT,
                            type="png",
                            animations="disabled",
                        )

                        # 验证截图文件
                        if not screenshot_output_path.exists() or screenshot_output_path.stat().st_size == 0:
                            return None, f"截图文件生成失败: {source_path}"

                        # 嵌入AIGC签名元数据
                        success = await self.embed_image_metadata(str(screenshot_output_path), aigc_params)
                        if success:
                            logger.info(f"PPT转换：成功嵌入签名元数据到图片 {screenshot_output_path}")
                        else:
                            logger.error(f"PPT转换：签名元数据嵌入失败 {screenshot_output_path}")
                            return None, f"签名元数据嵌入失败: {source_path}"

                        # 不再重命名截图文件，保持原始文件名

                        # 重置 agent idle 时间
                        self.update_agent_activity("PPTX转换中")

                        # 🎯 实时更新截图进度 - 截图阶段占总进度的90%
                        if task_mgr and task_key:
                            with progress_lock:
                                completed_screenshots += 1
                                # 🎯 修复进度计算：使用整个任务的总文件数量，而不是单个项目的文件数量
                                # 截图阶段占90%，计算当前截图进度
                                screenshot_progress = (completed_screenshots / total_valid_files) * 90

                            # 异步更新进度，避免阻塞
                            try:
                                await task_mgr.update_conversion_rate(
                                    task_key,
                                    min(screenshot_progress, 90.0),  # 确保不超过90%
                                    completed_screenshots,
                                    total_valid_files,  # 使用整个任务的总文件数量
                                )
                            except Exception as progress_error:
                                logger.warning(f"更新截图进度失败 {project_name}: {progress_error}")

                        logger.debug(f"项目 '{project_name}' 截图生成成功: {source_path} -> {screenshot_output_path}")
                        return screenshot_output_path, None

                    except Exception as page_error:
                        return None, f"页面处理失败 {source_path}: {str(page_error)}"
                    finally:
                        if page:
                            await page.close()

                except Exception as exc:
                    return None, f"文件转换异常 {source_path}: {str(exc)}"

            # 🎯 串行处理：逐个处理文件，避免网络资源竞争
            results = []

            logger.info(f"项目 '{project_name}' 开始截图转换，共 {len(files_to_convert)} 个文件")

            try:
                # 串行处理每个文件
                for i, file_path in enumerate(files_to_convert):
                    # 每处理5个文件输出一次进度，避免过多日志
                    if i % 5 == 0 or i == len(files_to_convert) - 1:
                        logger.info(f"项目 '{project_name}' 截图进度: {i + 1}/{len(files_to_convert)}")
                    result = await process_single_file_to_screenshot(file_path, i)
                    results.append(result)

            except Exception as e:
                logger.error(f"项目 '{project_name}' 串行任务执行失败: {e}")
                raise RuntimeError(f"项目 '{project_name}' 截图转换失败: {str(e)}")

            # 处理结果
            for i, result in enumerate(results):
                try:
                    if isinstance(result, (Exception, BaseException)):
                        logger.error(f"项目 '{project_name}' 处理文件时发生异常: {result}")
                        conversion_errors.append(
                            f"文件 {files_to_convert[i] if i < len(files_to_convert) else f'文件{i + 1}'}: 处理异常 - {str(result)}"
                        )
                        continue

                    screenshot_path, error_msg = result
                    if error_msg:
                        conversion_errors.append(error_msg)
                        logger.warning(f"项目 '{project_name}' 文件转换失败: {error_msg}")
                    elif screenshot_path:
                        screenshot_files.append(screenshot_path)

                except Exception as process_error:
                    logger.error(f"项目 '{project_name}' 处理转换结果时发生错误: {process_error}")
                    conversion_errors.append(
                        f"文件 {files_to_convert[i] if i < len(files_to_convert) else f'文件{i + 1}'}: 结果处理异常 - {str(process_error)}"
                    )

            logger.info(
                f"项目 '{project_name}' 串行截图转换完成: {len(screenshot_files)}/{len(files_to_convert)} 成功 (避免网络资源竞争)"
            )
            return screenshot_files, conversion_errors
        finally:
            await self._close_shared_browser_context(
                playwright_instance, shared_browser, shared_context, log_prefix=f"项目 '{project_name}'"
            )

    async def _convert_projects(
        self,
        projects: Dict[str, Dict[str, Any]],
        output_dir: Path,
        options: Optional[Dict[str, Any]] = None,
        task_mgr=None,
        task_key: Optional[str] = None,
        valid_files_count: int = 0,
        optimal_concurrency: int = 1,
        aigc_params: Optional[AigcMetadataParams] = None,
    ) -> tuple[List[Path], List[str]]:
        """
        实现抽象方法：执行PPTX项目转换

        Args:
            projects: PPTX项目字典
            output_dir: 输出目录（对应batch_dir）
            options: 转换选项（已忽略，不使用）
            task_mgr: 任务管理器
            task_key: 任务键
            valid_files_count: 有效文件数量
            optimal_concurrency: 最优并发数（用于计算截图并发数）
            aigc_params: AIGC元数据参数对象

        Returns:
            (转换后的PPTX文件列表, 错误信息列表)

        Note:
            此方法忽略options参数，使用系统默认的转换配置
        """
        # 注意：CDN 资源自动通过 Playwright 路由拦截映射到本地 static/
        generated_pptx_files = []
        all_conversion_errors = []
        total_processed_files = 0

        # 创建screenshots子目录
        screenshot_dir = output_dir / "screenshots"
        screenshot_dir.mkdir(parents=True, exist_ok=True)

        for project_name, project_info in projects.items():
            try:
                logger.info(f"开始处理PPTX项目: {project_name}")

                # 为截图操作计算专门的并发数（比普通转换更保守，避免内存爆炸）
                # 当前截图流程为串行处理以避免资源竞争，固定并发为1
                screenshot_concurrency = 1

                pptx_file, conversion_errors = await self._create_pptx_for_project(
                    project_name,
                    project_info,
                    screenshot_dir,
                    output_dir,
                    task_mgr,
                    task_key,
                    total_processed_files,
                    screenshot_concurrency,
                    valid_files_count,
                    aigc_params,
                )

                if pptx_file:
                    generated_pptx_files.append(pptx_file)
                    logger.info(f"PPTX项目 '{project_name}' 处理完成: {pptx_file}")

                all_conversion_errors.extend(conversion_errors)
                total_processed_files += len(project_info["files"])

            except Exception as e:
                logger.error(f"处理PPTX项目 '{project_name}' 失败: {e}")
                all_conversion_errors.append(f"项目 '{project_name}': {str(e)}")

        # 结束内存监控
        # self._end_memory_monitoring(task_key, len(generated_pptx_files), len(projects))  # 注释掉内存监控日志

        logger.info(f"所有PPTX项目处理完成: 成功 {len(generated_pptx_files)}/{len(projects)} 个项目")

        return generated_pptx_files, all_conversion_errors

    async def _get_service_specific_result_data(
        self, file_keys: List[Dict[str, str]], projects: Dict[str, Dict[str, Any]], converted_files: List[Path]
    ) -> Dict[str, Any]:
        """
        实现抽象方法：获取PPTX服务特定的结果数据

        Args:
            file_keys: 原始文件键列表
            projects: 项目字典
            converted_files: 转换后的PPTX文件列表

        Returns:
            PPTX服务特定的结果数据字典
        """
        return {
            "total_files": len(file_keys),
            "ppt_projects": len(projects),
            "generated_pptx_count": len(converted_files),
        }
